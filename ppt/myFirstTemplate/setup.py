from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
import plotly.graph_objects as go
import numpy as np
import pandas as pd
import datetime

images_path = 'images/'

def date_to_show(date):
    return datetime.datetime.strptime(date, '%Y-%m-%d').strftime('%d-%m-%Y')

def logo_path(number):
    if number == 1:
        return images_path + 'logo.png'
    if number == 2:
        return images_path + 'logo2.png'
    if number == 3:
        return images_path + 'logo3.png'

def pallete():
    colors = {'color_primary' : RGBColor(238, 238, 238), 'color_secondary': RGBColor(1, 12, 18), 'color_accent': RGBColor(115, 26, 4), 'color_accent_dark': RGBColor(29, 35, 46)}
    return colors

def presentation():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    return prs


def new_slide(prs,layout_number):
    return prs.slides.add_slide(prs.slide_layouts[layout_number])


def new_title(slide, title):
    slide.shapes.title.text = title
    return slide.shapes.title.text


def new_subtitle(slide, subtitle):
    slide.placeholders[1].text = subtitle
    return slide.placeholders[1].text

# print(pallete()['color_primary'])
