from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
import plotly.graph_objects as go
import numpy as np
import pandas as pd
import datetime
from setup import images_path, pallete, date_to_show, logo_path, presentation, new_slide, new_title, new_subtitle


# Setup Project
#------------------------------------------------------------------------------

file_name = input('Enter file name: ')
presentation_title = input('Enter presentation title: ')
colors = pallete()
presentation_subtitle = input('Enter presentation subtitle: ')
# last_worked_on = datetime.datetime.now().strftime('%d-%m-%Y')
# date_presented = date_to_show(input('Enter date presented 20-05-2021: '))
# author = input('Enter presentation author: ')
# company = input('Enter company name: ')


#------------------------------------------------------------------------------
# Front Page
#------------------------------------------------------------------------------
prs = presentation()
slide = new_slide(prs, layout_number=0)

new_title(slide, presentation_title)
new_subtitle(slide, presentation_subtitle)
# slide.shapes.title.text = presentation_title
# slide.placeholders[1].text = presentation_subtitle

slide.first_slide = True # Do I need this?
background_fill = slide.background.fill
background_fill.solid()
background_fill.fore_color.rgb = colors['color_accent']

#------------------------------------------------------------------------------
# Page 1
#------------------------------------------------------------------------------
slide = new_slide(prs, layout_number=1)




# Print the resulting powerpoint file
# ----------------------------------------------------------------------------------------------------------------------
prs.save(file_name + '.pptx')
