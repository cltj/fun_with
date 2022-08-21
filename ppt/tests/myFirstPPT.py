from pptx import presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE


# Create a presentation and add blank slide
prs = presentation.Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "This is text inside a textbox"

